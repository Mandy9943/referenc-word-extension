#!/usr/bin/env python3
"""Convert a PPTX into a DOCX with one slide image + speaker notes per slide.

Pipeline:
1) Export PPTX -> PDF (PowerPoint AppleScript on macOS, LibreOffice on Linux)
2) pdftoppm render: PDF -> PNG pages
3) python-pptx extract speaker notes
4) python-docx assemble final document
"""

from __future__ import annotations

import argparse
import re
import shutil
import subprocess
import sys
from pathlib import Path

try:
    from docx import Document
    from docx.shared import Inches
    from pptx import Presentation
    IMPORT_ERROR: Exception | None = None
except ModuleNotFoundError as exc:
    IMPORT_ERROR = exc

INVALID_XML_CHARS_RE = re.compile(r"[\x00-\x08\x0B\x0C\x0E-\x1F]")


def run_command(cmd: list[str], timeout_sec: int | None = None) -> None:
    try:
        proc = subprocess.run(
            cmd,
            text=True,
            capture_output=True,
            timeout=timeout_sec,
        )
    except subprocess.TimeoutExpired as exc:
        raise RuntimeError(
            f"Command timed out after {timeout_sec}s: {' '.join(cmd)}"
        ) from exc

    if proc.returncode != 0:
        raise RuntimeError(
            f"Command failed ({proc.returncode}): {' '.join(cmd)}\n"
            f"stdout:\n{proc.stdout}\n"
            f"stderr:\n{proc.stderr}"
        )


def export_pdf_with_powerpoint(pptx_path: Path, pdf_path: Path) -> None:
    applescript = r'''
on run argv
    set inPath to item 1 of argv
    set outPath to item 2 of argv
    set outFile to POSIX file outPath

    tell application "Microsoft PowerPoint"
        open POSIX file inPath
        set p to active presentation
        save p in outFile as save as PDF
        close p saving no
    end tell
end run
'''

    cmd = ["osascript", "-e", applescript, str(pptx_path), str(pdf_path)]
    try:
        run_command(cmd, timeout_sec=240)
    except RuntimeError:
        # Fallback for files that fail with direct AppleScript open/save but
        # can still be opened via Launch Services.
        run_command(["open", "-a", "Microsoft PowerPoint", str(pptx_path)], timeout_sec=60)
        fallback_script = r'''
on run argv
    set inPath to item 1 of argv
    set outPath to item 2 of argv
    set inFileName to do shell script "basename " & quoted form of inPath
    set outFile to POSIX file outPath
    set targetPres to missing value

    tell application "Microsoft PowerPoint"
        repeat with i from 1 to 60
            repeat with p in presentations
                try
                    if (name of p) is inFileName then
                        set targetPres to p
                        exit repeat
                    end if
                end try
            end repeat
            if targetPres is not missing value then exit repeat
            delay 1
        end repeat

        if targetPres is missing value then
            if (count of presentations) = 0 then
                error "PowerPoint did not open any presentation in fallback mode."
            end if
            set targetPres to active presentation
        end if

        save targetPres in outFile as save as PDF
        close targetPres saving no
    end tell
end run
'''
        run_command(
            ["osascript", "-e", fallback_script, str(pptx_path), str(pdf_path)],
            timeout_sec=240,
        )
    if not pdf_path.exists():
        raise RuntimeError("PowerPoint export completed but PDF was not created.")


def export_pdf_with_soffice(pptx_path: Path, pdf_path: Path) -> None:
    soffice_cmd = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice_cmd:
        raise RuntimeError(
            "LibreOffice CLI was not found (`soffice`/`libreoffice`). "
            "Install LibreOffice or run this script on macOS with Microsoft PowerPoint."
        )

    out_dir = pdf_path.parent
    out_dir.mkdir(parents=True, exist_ok=True)
    expected_pdf = out_dir / f"{pptx_path.stem}.pdf"
    if expected_pdf.exists():
        expected_pdf.unlink(missing_ok=True)

    run_command(
        [
            soffice_cmd,
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(out_dir),
            str(pptx_path),
        ],
        timeout_sec=240,
    )

    if not expected_pdf.exists():
        raise RuntimeError("LibreOffice export completed but PDF was not created.")

    if expected_pdf != pdf_path:
        pdf_path.unlink(missing_ok=True)
        shutil.move(str(expected_pdf), str(pdf_path))


def export_pdf(pptx_path: Path, pdf_path: Path) -> None:
    if sys.platform == "darwin" and shutil.which("osascript"):
        try:
            export_pdf_with_powerpoint(pptx_path, pdf_path)
            return
        except Exception:
            # If PowerPoint export fails, fall back to LibreOffice where available.
            if not (shutil.which("soffice") or shutil.which("libreoffice")):
                raise

    export_pdf_with_soffice(pptx_path, pdf_path)


def render_pdf_to_pngs(pdf_path: Path, output_prefix: Path, dpi: int) -> list[Path]:
    cmd = [
        "pdftoppm",
        "-png",
        "-r",
        str(dpi),
        str(pdf_path),
        str(output_prefix),
    ]
    run_command(cmd)

    png_files = sorted(
        output_prefix.parent.glob(f"{output_prefix.name}-*.png"),
        key=lambda p: int(re.search(r"-(\d+)\.png$", p.name).group(1)),
    )
    if not png_files:
        raise RuntimeError("No slide images were generated from the exported PDF.")
    return png_files


def extract_speaker_notes(pptx_path: Path) -> list[str]:
    prs = Presentation(str(pptx_path))
    notes: list[str] = []

    for slide in prs.slides:
        text = ""
        if slide.has_notes_slide:
            try:
                text = slide.notes_slide.notes_text_frame.text or ""
            except Exception:
                text = ""

        text = sanitize_for_docx(text.strip())
        if text == "Click to add notes":
            text = ""
        notes.append(text)

    return notes


def sanitize_for_docx(text: str) -> str:
    """Remove control characters disallowed by XML 1.0/WordprocessingML."""
    return INVALID_XML_CHARS_RE.sub("", text)


def make_safe_stem(stem: str) -> str:
    safe = re.sub(r"[^A-Za-z0-9._-]+", "_", stem).strip("._")
    return safe or "presentation"


def choose_work_dir(output_docx: Path) -> Path:
    """Pick a stable work directory for conversion artifacts."""
    home = Path.home()
    if sys.platform == "darwin":
        candidates = [
            home / "Library/Group Containers/UBF8T346G9.Office/pptx_to_docx_work",
            home / "Library/Containers/com.microsoft.Powerpoint/Data/pptx_to_docx_work",
            output_docx.parent / ".pptx_to_docx_work",
        ]
    else:
        candidates = [output_docx.parent / ".pptx_to_docx_work", home / ".pptx_to_docx_work"]

    for candidate in candidates:
        try:
            candidate.mkdir(parents=True, exist_ok=True)
            return candidate
        except Exception:
            continue

    raise RuntimeError(
        "Unable to create a writable work directory for conversion."
    )


def prepare_work_paths(
    input_pptx: Path, output_docx: Path
) -> tuple[Path, Path, Path, str]:
    """Use a stable on-disk work folder (no random /private/var temp dirs)."""
    safe_stem = make_safe_stem(input_pptx.stem)
    work_dir = choose_work_dir(output_docx)
    pdf_path = work_dir / f"{safe_stem}.pdf"
    image_prefix = work_dir / f"{safe_stem}_slide"
    return work_dir, pdf_path, image_prefix, safe_stem


def cleanup_run_artifacts(work_dir: Path, safe_stem: str) -> None:
    for png in work_dir.glob(f"{safe_stem}_slide-*.png"):
        png.unlink(missing_ok=True)
    (work_dir / f"{safe_stem}.pdf").unlink(missing_ok=True)
    (work_dir / f"{safe_stem}__input.pptx").unlink(missing_ok=True)


def stage_input_for_conversion(input_pptx: Path, work_dir: Path, safe_stem: str) -> Path:
    """Copy source deck into a stable work folder for conversion tools."""
    staged = work_dir / f"{safe_stem}__input.pptx"
    shutil.copy2(input_pptx, staged)
    return staged


def build_docx(slide_images: list[Path], notes: list[str], output_docx: Path) -> None:
    if len(slide_images) != len(notes):
        raise RuntimeError(
            f"Slide/image mismatch: {len(slide_images)} images vs {len(notes)} slides"
        )

    doc = Document()
    section = doc.sections[0]

    max_width_inches = (
        section.page_width - section.left_margin - section.right_margin
    ) / 914400

    for index, (image_path, note_text) in enumerate(zip(slide_images, notes), start=1):
        doc.add_heading(f"Slide {index}", level=2)
        doc.add_picture(str(image_path), width=Inches(max_width_inches))

        label = doc.add_paragraph()
        label.add_run("Speaker Notes").bold = True

        body = doc.add_paragraph()
        body.text = note_text if note_text else "[No speaker notes]"

        if index != len(slide_images):
            doc.add_page_break()

    output_docx.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(output_docx))


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Convert PPTX to DOCX with each slide image and speaker notes."
    )
    parser.add_argument(
        "input_pptx",
        nargs="*",
        type=Path,
        help="One or more source .pptx files",
    )
    parser.add_argument(
        "-o",
        "--output",
        type=Path,
        default=None,
        help="Output .docx path (default: <pptx_stem>_slides_notes.docx)",
    )
    parser.add_argument(
        "--batch-desktop",
        action="store_true",
        help="Batch convert all .pptx files on Desktop.",
    )
    parser.add_argument(
        "--latest-desktop",
        action="store_true",
        help="Convert the most recently modified .pptx on Desktop.",
    )
    parser.add_argument(
        "--dpi",
        type=int,
        default=180,
        help="Image render DPI for slide PNGs (default: 180)",
    )
    parser.add_argument(
        "--keep-temp",
        action="store_true",
        help="Keep temporary PDF/PNG files for inspection",
    )
    return parser.parse_args()


def resolve_inputs(args: argparse.Namespace) -> list[Path]:
    selected: list[Path] = []

    for path in args.input_pptx:
        resolved = path.expanduser().resolve()
        if resolved not in selected:
            selected.append(resolved)

    desktop = Path.home() / "Desktop"
    if args.batch_desktop:
        for path in sorted(desktop.glob("*.pptx")):
            resolved = path.resolve()
            if resolved not in selected:
                selected.append(resolved)

    if args.latest_desktop:
        desktop_pptx = [p for p in desktop.glob("*.pptx") if p.is_file()]
        if not desktop_pptx:
            raise RuntimeError("No .pptx files found on Desktop for --latest-desktop.")
        latest = max(desktop_pptx, key=lambda p: p.stat().st_mtime).resolve()
        if latest not in selected:
            selected.append(latest)

    return selected


def convert_one(
    pptx_path: Path,
    output_docx: Path,
    dpi: int,
    keep_temp: bool,
) -> int:
    if not pptx_path.exists() or pptx_path.suffix.lower() != ".pptx":
        raise RuntimeError(f"Input must be an existing .pptx file: {pptx_path}")

    work_dir, pdf_path, image_prefix, safe_stem = prepare_work_paths(pptx_path, output_docx)

    cleanup_run_artifacts(work_dir, safe_stem)
    staged_pptx = stage_input_for_conversion(pptx_path, work_dir, safe_stem)
    export_pdf(staged_pptx, pdf_path)
    slide_images = render_pdf_to_pngs(pdf_path, image_prefix, dpi)
    notes = extract_speaker_notes(staged_pptx)
    build_docx(slide_images, notes, output_docx)

    if keep_temp:
        kept_path = output_docx.with_suffix("").with_name(output_docx.stem + "_tmp")
        if kept_path.exists():
            shutil.rmtree(kept_path)
        kept_path.mkdir(parents=True, exist_ok=True)
        shutil.copy2(pdf_path, kept_path / pdf_path.name)
        for png in work_dir.glob(f"{safe_stem}_slide-*.png"):
            shutil.copy2(png, kept_path / png.name)
        print(f"Kept temp files at: {kept_path}")
    else:
        cleanup_run_artifacts(work_dir, safe_stem)

    return len(slide_images)


def main() -> int:
    if IMPORT_ERROR is not None:
        print(
            "Missing dependency. Run with uv to auto-provision packages:",
            file=sys.stderr,
        )
        print(
            f"  uv run --with python-docx --with python-pptx {Path(__file__)} <input.pptx>",
            file=sys.stderr,
        )
        print(f"Underlying error: {IMPORT_ERROR}", file=sys.stderr)
        return 3

    args = parse_args()
    try:
        inputs = resolve_inputs(args)
    except Exception as exc:
        print(f"Error: {exc}", file=sys.stderr)
        return 1

    if not inputs:
        print(
            "No input .pptx files provided. Use file paths, --latest-desktop, or --batch-desktop.",
            file=sys.stderr,
        )
        return 2

    if args.output is not None and len(inputs) != 1:
        print("--output can only be used when converting a single input file.", file=sys.stderr)
        return 2

    failures = 0
    for pptx_path in inputs:
        output_docx = (
            args.output.expanduser().resolve()
            if args.output is not None
            else pptx_path.with_name(f"{pptx_path.stem}_slides_notes.docx")
        )
        try:
            slide_count = convert_one(
                pptx_path=pptx_path,
                output_docx=output_docx,
                dpi=args.dpi,
                keep_temp=args.keep_temp,
            )
            print(f"Created: {output_docx}")
            print(f"Slides processed: {slide_count}")
        except Exception as exc:
            failures += 1
            print(f"Error converting {pptx_path}: {exc}", file=sys.stderr)

    if len(inputs) > 1:
        print(f"Batch summary: {len(inputs) - failures}/{len(inputs)} succeeded.")

    return 0 if failures == 0 else 1


def cli() -> None:
    raise SystemExit(main())


if __name__ == "__main__":
    cli()
